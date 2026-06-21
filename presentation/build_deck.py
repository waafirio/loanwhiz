"""Generate the LoanWhiz hackathon presentation (.pptx).

Clean, professional, light-theme B2B-SaaS design system. No external assets;
everything is drawn with python-pptx shapes so it renders identically anywhere.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Design tokens
# ---------------------------------------------------------------------------
INK      = RGBColor(0x16, 0x20, 0x33)   # near-black navy (primary text)
NAVY     = RGBColor(0x1B, 0x2A, 0x4A)   # deep navy (headers / bars)
BLUE     = RGBColor(0x2B, 0x59, 0xD6)   # primary brand blue
TEAL     = RGBColor(0x12, 0x9E, 0x8F)   # accent teal
SLATE    = RGBColor(0x53, 0x61, 0x77)   # secondary text
MIST     = RGBColor(0x6B, 0x7A, 0x90)   # tertiary / captions
CLOUD    = RGBColor(0xF4, 0xF7, 0xFB)   # card background
LINE     = RGBColor(0xDD, 0xE4, 0xEE)   # hairline borders
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
AMBER    = RGBColor(0xB7, 0x79, 0x10)   # caution / flag
GREEN    = RGBColor(0x1E, 0x8E, 0x4E)   # success

FONT   = "Calibri"
FONT_H = "Calibri"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, shadow=False, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if round_:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow:
        _soft_shadow(shp)
    return shp


def _soft_shadow(shp):
    spPr = shp._element.spPr
    el = spPr.makeelement(qn('a:effectLst'), {})
    sh = el.makeelement(qn('a:outerShdw'),
                        {'blurRad': '90000', 'dist': '40000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = sh.makeelement(qn('a:srgbClr'), {'val': '1B2A4A'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '16000'})
    clr.append(alpha); sh.append(clr); el.append(sh); spPr.append(el)


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, *rest) in para:
            italic = rest[0] if rest else False
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic
            r.font.name = FONT
    return tb


def page_chrome(s, n, title=None, kicker=None):
    """Standard content-slide chrome: accent bar, title, footer."""
    rect(s, 0, 0, EMU_W, EMU_H, fill=WHITE)
    # left accent rail
    rect(s, 0, 0, Inches(0.18), EMU_H, fill=BLUE)
    if title:
        if kicker:
            text(s, Inches(0.7), Inches(0.42), Inches(11.8), Inches(0.3),
                 [[(kicker.upper(), 12, TEAL, True)]])
        text(s, Inches(0.7), Inches(0.72), Inches(11.9), Inches(0.9),
             [[(title, 28, NAVY, True)]])
        rect(s, Inches(0.72), Inches(1.46), Inches(0.6), Pt(3), fill=TEAL)
    # footer
    text(s, Inches(0.7), Inches(7.04), Inches(6), Inches(0.3),
         [[("LoanWhiz", 9, SLATE, True), ("  ·  Structured Finance Agent Framework", 9, MIST, False)]])
    text(s, Inches(11.8), Inches(7.04), Inches(0.9), Inches(0.3),
         [[(f"{n:02d}", 9, MIST, False)]], align=PP_ALIGN.RIGHT)


def chip(s, x, y, label, fg, bg, w=None):
    w = w or Inches(1.6)
    c = rect(s, x, y, w, Inches(0.34), fill=bg, round_=True)
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = fg; r.font.name = FONT
    return c


def card(s, x, y, w, h, title, body_runs, accent=BLUE, title_size=15, num=None):
    rect(s, x, y, w, h, fill=CLOUD, line=LINE, line_w=1, shadow=True, round_=True)
    rect(s, x, y, Inches(0.09), h, fill=accent, round_=False)
    ty = y + Inches(0.22)
    if num is not None:
        text(s, x + Inches(0.28), ty, Inches(0.7), Inches(0.4),
             [[(num, 13, accent, True)]])
    text(s, x + Inches(0.28), ty, w - Inches(0.5), Inches(0.4),
         [[(title, title_size, NAVY, True)]])
    text(s, x + Inches(0.28), ty + Inches(0.46), w - Inches(0.55), h - Inches(0.7),
         body_runs, line_spacing=1.06, space_after=3)


def bullets(items, size=13.5, color=SLATE, gap=6, head_color=INK):
    """Build paragraph runs with a teal tick prefix."""
    out = []
    for it in items:
        if isinstance(it, tuple):
            head, rest = it
            out.append([("▸  ", size, TEAL, True), (head, size, head_color, True),
                        (rest, size, color, False)])
        else:
            out.append([("▸  ", size, TEAL, True), (it, size, color, False)])
    return out


# ===========================================================================
# 1 — TITLE
# ===========================================================================
s = slide()
rect(s, 0, 0, EMU_W, EMU_H, fill=NAVY)
# subtle accent band
rect(s, 0, Inches(5.55), EMU_W, Inches(0.06), fill=TEAL)
rect(s, 0, 0, EMU_W, Inches(0.06), fill=TEAL)
text(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.4),
     [[("BARCELONA AI TINKERERS  ·  STRUCTURED FINANCE HACKATHON 2026", 13, RGBColor(0x8F,0xB6,0xF0), True)]])
text(s, Inches(0.86), Inches(2.15), Inches(11.6), Inches(1.4),
     [[("LoanWhiz", 66, WHITE, True)]])
text(s, Inches(0.9), Inches(3.35), Inches(11.2), Inches(1.0),
     [[("A Structured Finance ", 28, RGBColor(0xCF,0xDB,0xEE), False),
       ("Agent Framework", 28, TEAL, True)]])
text(s, Inches(0.92), Inches(4.35), Inches(11.0), Inches(0.9),
     [[("Turns a 300-page prospectus and ESMA loan tapes into a machine-executable,", 15, RGBColor(0xB9,0xC6,0xDC), False)],
      [("fully-audited deal model that an agent can reason over and operate.", 15, RGBColor(0xB9,0xC6,0xDC), False)]],
     line_spacing=1.15)
text(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.9),
     [[("Challenge 1 — Agent Framework", 13, WHITE, True),
       ("      Gemini 2.5 (Vertex AI) · LangGraph · Docling · FINOS Governance", 13, RGBColor(0x9F,0xB0,0xCC), False)]])
text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.4),
     [[("github.com/waafirio/loanwhiz", 12, RGBColor(0x7E,0xA6,0xE8), False)]])

# ===========================================================================
# 2 — THE PROBLEM
# ===========================================================================
s = slide(); page_chrome(s, 2, "The deal logic is locked inside documents", kicker="The problem")
text(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.7),
     [[("Every securitisation encodes its rules — ", 15, SLATE, False),
       ("who gets paid, in what order, under which triggers", 15, INK, True),
       (" — in prose, not code.", 15, SLATE, False)]], line_spacing=1.1)
cw, ch, gap = Inches(3.83), Inches(3.5), Inches(0.2)
x0 = Inches(0.7); y0 = Inches(2.7)
card(s, x0, y0, cw, ch, "Prospectus", bullets([
    ("~300 pages. ", "Waterfalls, triggers, covenants and definitions written as legal English."),
    ("Cross-referenced. ", "A single payment step pulls in a dozen defined terms."),
    ("Scanned / OCR-heavy. ", "Tables and clauses that break naive parsers."),
]), accent=BLUE)
card(s, x0+cw+gap, y0, cw, ch, "ESMA loan tapes", bullets([
    ("Loan-level. ", "Annex-2 RMBS tapes, tens of thousands of rows per period."),
    ("Monthly. ", "Performance, arrears and collections that move every reporting date."),
    ("Must reconcile. ", "Against the prospectus rules and the investor report."),
]), accent=TEAL)
card(s, x0+2*(cw+gap), y0, cw, ch, "Today it's manual", bullets([
    ("Analysts read by hand. ", "Slow, expensive, hard to audit, easy to get wrong."),
    ("Chatbots can't execute. ", "Naive RAG retrieves text; it can't run a waterfall."),
    ("No provenance. ", "An answer with no citation is useless to a trustee."),
]), accent=NAVY)
text(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.4),
     [[("The hackathon asks for an AI solution finance can ", 13.5, SLATE, False),
       ("trust and operate", 13.5, NAVY, True),
       (" — not a document chatbot.", 13.5, SLATE, False)]])

# ===========================================================================
# 3 — THE INSIGHT / APPROACH
# ===========================================================================
s = slide(); page_chrome(s, 3, "Don't chat about the deal — compile it", kicker="Our approach")
text(s, Inches(0.7), Inches(1.75), Inches(11.9), Inches(0.7),
     [[("We separate three concerns that a chatbot collapses into one. ", 15, SLATE, False),
       ("Read once, execute deterministically, orchestrate with an agent.", 15, INK, True)]],
     line_spacing=1.1)
cw, ch, gap = Inches(3.83), Inches(3.45), Inches(0.2)
x0 = Inches(0.7); y0 = Inches(2.7)
card(s, x0, y0, cw, ch, "1 · Extract", bullets([
    ("Prospectus → deal model. ", "Docling + section-routed Gemini 2.5 Pro produce structured JSON."),
    ("Done once, cached. ", "Extraction is the expensive step; never repeated at query time."),
    ("Every fact cited. ", "Verbatim excerpt + section for each extracted value."),
]), accent=BLUE, title_size=17)
card(s, x0+cw+gap, y0, cw, ch, "2 · Execute", bullets([
    ("SF-native primitives. ", "Waterfall, covenant monitor, projector — typed, testable, deterministic."),
    ("No LLM in the maths. ", "The agent calls code; numbers are computed, not generated."),
    ("Composable. ", "A registry the agent and UI both discover."),
]), accent=TEAL, title_size=17)
card(s, x0+2*(cw+gap), y0, cw, ch, "3 · Orchestrate", bullets([
    ("LangGraph agent. ", "Plans which primitives to call, in what order, then explains."),
    ("Grounded answers. ", "Responses cite the deal model and tape rows they used."),
    ("Governed. ", "Every call carries an audit entry and confidence score."),
]), accent=NAVY, title_size=17)
text(s, Inches(0.7), Inches(6.52), Inches(11.9), Inches(0.5),
     [[("Key insight:  ", 14, TEAL, True),
       ("the LLM's job is comprehension and orchestration — never arithmetic. That's what makes it auditable.", 14, INK, True)]])

# ===========================================================================
# 4 — ARCHITECTURE
# ===========================================================================
s = slide(); page_chrome(s, 4, "Architecture", kicker="How it fits together")

def box(x, y, w, h, title, sub, fill, tcol, sub_col=None, tsize=13):
    rect(s, x, y, w, h, fill=fill, line=LINE, line_w=1, round_=True, shadow=True)
    text(s, x+Inches(0.16), y+Inches(0.1), w-Inches(0.3), h-Inches(0.2),
         [[(title, tsize, tcol, True)]] + ([[ (sub, 10.5, sub_col or SLATE, False) ]] if sub else []),
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02, space_after=2, align=PP_ALIGN.CENTER)

def arrow(x, y, w, h, vert=False):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW if vert else MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xAF,0xBD,0xD2)
    a.line.fill.background(); a.shadow.inherit = False

top = Inches(1.7)
# Clients row
box(Inches(0.7), top, Inches(5.55), Inches(0.85), "Demo UI  ·  Next.js dashboard + docked chat", "5 deal views — Overview · Pool · Waterfall · Compliance · Projection", WHITE, NAVY)
box(Inches(6.95), top, Inches(5.65), Inches(0.85), "REST API  ·  FastAPI", "/deals · /deal/{id}/model · /waterfall · /compliance · /primitives · /query", WHITE, NAVY)
arrow(Inches(6.45), top+Inches(0.95), Inches(0.55), Inches(0.45), vert=True)

# Agent layer
ay = Inches(2.95)
rect(s, Inches(0.7), ay, Inches(11.9), Inches(1.15), fill=RGBColor(0xEC,0xF1,0xFB), line=BLUE, line_w=1.25, round_=True, shadow=True)
text(s, Inches(0.95), ay+Inches(0.12), Inches(11.4), Inches(0.4),
     [[("LANGGRAPH AGENT SERVICE", 12.5, BLUE, True),
       ("    Planner → DAG executor → validator → confidence scorer → human-review router → audit trail", 12, SLATE, False)]])
text(s, Inches(0.95), ay+Inches(0.62), Inches(11.4), Inches(0.4),
     [[("Governance layer:  ", 11.5, TEAL, True),
       ("FINOS AI Governance Framework — citations, replay, model-risk classification on every call", 11.5, SLATE, False)]])
arrow(Inches(3.3), ay+Inches(1.25), Inches(0.5), Inches(0.4), vert=True)
arrow(Inches(9.4), ay+Inches(1.25), Inches(0.5), Inches(0.4), vert=True)

# Primitives + deal model
my = Inches(4.62)
box(Inches(0.7), my, Inches(5.75), Inches(1.18),
    "SF PRIMITIVES  (registry)",
    "tape_normaliser · collections_aggregator · waterfall_runner\ncovenant_monitor · report_verifier · audit_logger",
    CLOUD, NAVY)
box(Inches(6.85), my, Inches(5.75), Inches(1.18),
    "DEAL MODEL  (JSON, per deal)",
    "definitions{} · waterfalls[] (Revenue/Redemption/Post-Enf) · triggers[] · tranches[]\nextracted from the prospectus, cited, cached",
    CLOUD, NAVY)
arrow(Inches(6.05), my+Inches(1.28), Inches(0.5), Inches(0.35), vert=True)

# Data layer
dy = Inches(6.18)
box(Inches(0.7), dy, Inches(11.9), Inches(0.7),
    "DATA LAYER",
    "Docling extraction (prospectus → JSON)   ·   direct-read ESMA tape ingestion (CSV/parquet)   ·   HuggingFace: green-lion-2026 + green-lion-2024-2025 (27-month history)",
    NAVY, WHITE, sub_col=RGBColor(0xC7,0xD3,0xE8))

# ===========================================================================
# 5 — EXTRACTION PIPELINE
# ===========================================================================
s = slide(); page_chrome(s, 5, "From prospectus to executable model", kicker="The extraction pipeline")
steps = [
    ("PDF", "300-page\nprospectus", NAVY),
    ("Docling", "structure-aware\nPDF → markdown", BLUE),
    ("Section router", "scope each call to\nthe right section", TEAL),
    ("Gemini 2.5 Pro", "structured extraction\n(forced JSON)", BLUE),
    ("Deal model", "waterfalls · triggers\ntranches · definitions", NAVY),
]
n = len(steps); bw = Inches(2.06); gap = Inches(0.30)
x = Inches(0.7); y = Inches(2.0)
for i, (t, sub, col) in enumerate(steps):
    rect(s, x, y, bw, Inches(1.5), fill=CLOUD, line=LINE, line_w=1, round_=True, shadow=True)
    rect(s, x, y, bw, Inches(0.1), fill=col, round_=False)
    text(s, x+Inches(0.12), y+Inches(0.22), bw-Inches(0.24), Inches(0.5),
         [[(t, 14, NAVY, True)]], align=PP_ALIGN.CENTER)
    text(s, x+Inches(0.1), y+Inches(0.72), bw-Inches(0.2), Inches(0.7),
         [[(line, 10.5, SLATE, False)] for line in sub.split("\n")], align=PP_ALIGN.CENTER, line_spacing=1.02)
    if i < n-1:
        a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x+bw+Inches(0.02), y+Inches(0.58), gap-Inches(0.04), Inches(0.34))
        a.fill.solid(); a.fill.fore_color.rgb = RGBColor(0xAF,0xBD,0xD2); a.line.fill.background(); a.shadow.inherit=False
    x = x + bw + gap

# why not RAG
yy = Inches(3.95)
cw, gap2 = Inches(3.83), Inches(0.2)
card(s, Inches(0.7), yy, cw, Inches(2.45), "Why not naive RAG?", bullets([
    "Chunk-and-retrieve loses the ordered structure a waterfall depends on.",
    "It returns passages — you still can't execute the deal.",
    "Section-routing keeps each LLM call small, scoped and low-hallucination.",
], size=12.5), accent=AMBER)
card(s, Inches(0.7)+cw+gap2, yy, cw, Inches(2.45), "Zero-shot, no training", bullets([
    "No fine-tuning, no curated corpus — Gemini 2.5 Pro's pre-trained capability.",
    "Forced function-calling guarantees typed, schema-valid output.",
    "Portable to any ESMA-style RMBS/ABS prospectus.",
], size=12.5), accent=TEAL)
card(s, Inches(0.7)+2*(cw+gap2), yy, cw, Inches(2.45), "Cited & cached", bullets([
    "Every step, term and threshold carries a verbatim source excerpt.",
    "Confidence is a real coverage metric — sections resolved / steps with a recipient.",
    "Extracted once, cached durably under data/ — query time is instant.",
], size=12.5), accent=BLUE)

# ===========================================================================
# 6 — PRIMITIVES
# ===========================================================================
s = slide(); page_chrome(s, 6, "SF-native primitives", kicker="The framework's building blocks")
text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.55),
     [[("A registry of composable, typed operations. Each is independently testable, returns a confidence score and citations, ", 13.5, SLATE, False),
       ("and is discoverable by the agent and the UI.", 13.5, INK, True)]], line_spacing=1.08)
prims = [
    ("esma_tape_normaliser", "Normalise Annex-2 tapes → pool analytics: WAL, arrears, EPC / geo / rate distributions."),
    ("collections_aggregator", "Aggregate monthly interest, principal, prepayments and recoveries into waterfall inputs."),
    ("waterfall_runner", "Execute the extracted waterfall against a period's collections; per-tranche distributions + trace."),
    ("covenant_monitor", "Check tape metrics against extracted trigger thresholds; track breach proximity over time."),
    ("report_verifier", "Reconcile reconstructed collateral against the investor reports to the cent; flag discrepancies."),
    ("audit_logger", "Wrap every call with FINOS provenance: input/output hashes, confidence, citations, model version."),
]
cols, rows = 2, 3
cw, chh = Inches(5.9), Inches(1.02)
gx, gy = Inches(0.1), Inches(0.12)
x0, y0 = Inches(0.7), Inches(2.42)
for i, (name, desc) in enumerate(prims):
    r, c = i % rows, i // rows
    x = x0 + c*(cw+gx); y = y0 + r*(chh+gy)
    rect(s, x, y, cw, chh, fill=CLOUD, line=LINE, line_w=1, round_=True)
    rect(s, x, y, Inches(0.08), chh, fill=TEAL if c==0 else BLUE)
    text(s, x+Inches(0.24), y+Inches(0.12), cw-Inches(0.4), Inches(0.35),
         [[(name+"  ", 13, NAVY, True), ("v0.1.0", 9.5, MIST, False)]])
    text(s, x+Inches(0.24), y+Inches(0.45), cw-Inches(0.42), Inches(0.5),
         [[(desc, 11, SLATE, False)]], line_spacing=1.0)

# ===========================================================================
# 7 — GOVERNANCE
# ===========================================================================
s = slide(); page_chrome(s, 7, "Governance is baked into the interface", kicker="Trust — FINOS AI Governance Framework")
text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.55),
     [[("Not a post-hoc add-on. ", 14, INK, True),
       ("No primitive may return an output without a corresponding audit entry — the thing that makes this usable in finance.", 14, SLATE, False)]],
     line_spacing=1.08)
items = [
    ("Audit trail", "Append-only AuditEntry per call: input/output hashes, model version, timestamp, duration."),
    ("Confidence scoring", "Real coverage metric per primitive; agent answers aggregate it as min(per-tool confidence)."),
    ("Verbatim citations", "Every fact grounded to a section + verbatim excerpt; a reviewer can locate the source."),
    ("Replayability", "Input hash + model version identify each request; runs are comparable for drift."),
    ("Human-review routing", "Confidence < 0.70 → human_review_required = True; output can't be used until signed off."),
    ("Model-risk class", "Decision-support tier, zero autonomy — outputs are analytical inputs, never actions."),
]
cw, chh = Inches(3.83), Inches(1.62)
gx, gy = Inches(0.2), Inches(0.18)
x0, y0 = Inches(0.7), Inches(2.45)
for i, (t, d) in enumerate(items):
    r, c = i // 3, i % 3
    x = x0 + c*(cw+gx); y = y0 + r*(chh+gy)
    rect(s, x, y, cw, chh, fill=CLOUD, line=LINE, line_w=1, round_=True, shadow=True)
    text(s, x+Inches(0.22), y+Inches(0.16), cw-Inches(0.4), Inches(0.35), [[(t, 14, NAVY, True)]])
    rect(s, x+Inches(0.24), y+Inches(0.52), Inches(0.45), Pt(2.5), fill=TEAL)
    text(s, x+Inches(0.22), y+Inches(0.64), cw-Inches(0.42), Inches(0.9),
         [[(d, 11.5, SLATE, False)]], line_spacing=1.05)

# ===========================================================================
# 8 — WHAT WE BUILT / DEMO
# ===========================================================================
s = slide(); page_chrome(s, 8, "A working product, end-to-end", kicker="What we built")
text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5),
     [[("One loaded deal shared across five dashboard views plus a docked chat panel — all served over the same REST API.", 13.5, SLATE, False)]])
views = [
    ("Overview", "Extracted deal model — tranche structure, trigger names, completeness score."),
    ("Pool & Performance", "3-period pool analytics; arrears, EPC and geographic distributions."),
    ("Waterfall", "Revenue priority cascade and per-tranche distributions for the period."),
    ("Compliance", "Live covenant monitor across reporting periods, with breach proximity."),
    ("Projection", "Base vs stress forward cashflows, including Class-A WAL."),
    ("Docked chat", "Ad-hoc deal questions, grounded in the loaded deal model and tapes."),
]
cw, chh = Inches(3.83), Inches(1.42)
gx, gy = Inches(0.2), Inches(0.18)
x0, y0 = Inches(0.7), Inches(2.4)
for i, (t, d) in enumerate(views):
    r, c = i // 3, i % 3
    x = x0 + c*(cw+gx); y = y0 + r*(chh+gy)
    accent = TEAL if i == 5 else BLUE
    rect(s, x, y, cw, chh, fill=WHITE, line=LINE, line_w=1, round_=True, shadow=True)
    rect(s, x, y, cw, Inches(0.1), fill=accent)
    text(s, x+Inches(0.22), y+Inches(0.24), cw-Inches(0.4), Inches(0.35), [[(t, 14.5, NAVY, True)]])
    text(s, x+Inches(0.22), y+Inches(0.66), cw-Inches(0.42), Inches(0.7),
         [[(d, 11.5, SLATE, False)]], line_spacing=1.05)
text(s, Inches(0.7), Inches(6.62), Inches(11.9), Inches(0.4),
     [[("Surfaces the framework primitives directly in the UI — the challenge asks to see the building blocks, not just answers.", 12.5, MIST, False, True)]])

# ===========================================================================
# 9 — PROOF ON GREEN LION
# ===========================================================================
s = slide(); page_chrome(s, 9, "Demonstrated on Green Lion 2026-1", kicker="Results")
text(s, Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5),
     [[("A synthetic Dutch RMBS released for the hackathon: prospectus + 27 monthly ESMA tape snapshots (2024–2026) + 3 investor reports.", 13.5, SLATE, False)]])
# metric tiles
mets = [
    ("27", "monthly tapes", "Synthetic snapshots, Jan 2024 – Apr 2026 (Jan-2026 gap)"),
    ("11", "revenue steps", "Full Revenue Priority of Payments (a)–(k), cited"),
    ("3", "waterfalls", "Revenue · Redemption · Post-Enforcement"),
    ("3", "live triggers", "Class A & B PDL + reserve-fund shortfall"),
]
tw, gx = Inches(2.86), Inches(0.21)
x0, y0 = Inches(0.7), Inches(2.45)
for i, (big, lab, sub) in enumerate(mets):
    x = x0 + i*(tw+gx)
    rect(s, x, y0, tw, Inches(1.85), fill=CLOUD, line=LINE, line_w=1, round_=True, shadow=True)
    text(s, x, y0+Inches(0.18), tw, Inches(0.8), [[(big, 44, BLUE, True)]], align=PP_ALIGN.CENTER)
    text(s, x, y0+Inches(1.02), tw, Inches(0.35), [[(lab, 14, NAVY, True)]], align=PP_ALIGN.CENTER)
    text(s, x+Inches(0.18), y0+Inches(1.36), tw-Inches(0.36), Inches(0.45),
         [[(sub, 10, SLATE, False)]], align=PP_ALIGN.CENTER, line_spacing=1.0)
# bottom proof bar
by = Inches(4.65)
card(s, Inches(0.7), by, Inches(5.95), Inches(1.9), "End-to-end, verified", bullets([
    "Prospectus → extracted model → model-driven waterfall → multi-period state.",
    "Collateral reconciles to the investor reports to the cent; liabilities reconstructed from the prospectus + invariant-validated.",
    "Chat answers grounded numbers (e.g. pool balance ~€1bn) live.",
], size=12), accent=GREEN)
card(s, Inches(6.85), by, Inches(5.75), Inches(1.9), "Data-driven by design", bullets([
    "Add a deal via data/deals.json — no code change; the interpreter executes its extracted model.",
    "Format-agnostic tapes: CSV or parquet (incl. a combined multi-month file).",
    "Validated on Green Lion 2026-1; real multi-deal validation underway (#206).",
], size=12), accent=TEAL)

# ===========================================================================
# 10 — STACK + CLOSE
# ===========================================================================
s = slide(); page_chrome(s, 10, "Built on open foundations", kicker="Stack & what's next")
stack = [
    ("Gemini 2.5 Pro / Flash", "Extraction & agent orchestration (Vertex AI)"),
    ("LangGraph", "ReAct planner + DAG execution"),
    ("Docling (IBM)", "Structure-aware PDF → markdown"),
    ("FINOS AI Gov. Framework", "Audit, confidence, model-risk patterns"),
    ("deeploans (Algoritmica)", "ESMA loan-tape ingestion ETL"),
    ("FastAPI + Next.js", "REST API and the demo dashboard"),
]
cw, chh = Inches(3.83), Inches(0.95)
gx, gy = Inches(0.2), Inches(0.16)
x0, y0 = Inches(0.7), Inches(1.95)
for i, (t, d) in enumerate(stack):
    r, c = i // 3, i % 3
    x = x0 + c*(cw+gx); y = y0 + r*(chh+gy)
    rect(s, x, y, cw, chh, fill=CLOUD, line=LINE, line_w=1, round_=True)
    rect(s, x, y, Inches(0.08), chh, fill=BLUE)
    text(s, x+Inches(0.24), y+Inches(0.13), cw-Inches(0.4), Inches(0.35), [[(t, 13, NAVY, True)]])
    text(s, x+Inches(0.24), y+Inches(0.46), cw-Inches(0.42), Inches(0.4), [[(d, 11, SLATE, False)]])

# closing band
cy = Inches(4.15)
rect(s, Inches(0.7), cy, Inches(11.9), Inches(2.45), fill=NAVY, round_=True, shadow=True)
rect(s, Inches(0.7), cy, Inches(11.9), Inches(0.09), fill=TEAL)
text(s, Inches(1.05), cy+Inches(0.28), Inches(11.2), Inches(0.5),
     [[("Why it answers the challenge", 18, WHITE, True)]])
text(s, Inches(1.05), cy+Inches(0.85), Inches(11.2), Inches(1.4),
     [[("▸  ", 13.5, TEAL, True), ("A reusable agent framework", 13.5, WHITE, True),
       (" — SF-native primitives + a registry, not a one-off script.", 13.5, RGBColor(0xC7,0xD3,0xE8), False)],
      [("▸  ", 13.5, TEAL, True), ("Executes deal logic", 13.5, WHITE, True),
       (" — extracts a machine-runnable model and computes real numbers, deterministically.", 13.5, RGBColor(0xC7,0xD3,0xE8), False)],
      [("▸  ", 13.5, TEAL, True), ("Trustworthy by construction", 13.5, WHITE, True),
       (" — citations, confidence and human-review routing on every output.", 13.5, RGBColor(0xC7,0xD3,0xE8), False)],
      [("▸  ", 13.5, TEAL, True), ("Data-driven & extensible", 13.5, WHITE, True),
       (" — Apache-2.0; add a deal via data/deals.json (no code change), CSV or parquet tapes, add a primitive in a few lines. Demonstrated on Green Lion 2026-1; multi-deal validation underway.", 13.5, RGBColor(0xC7,0xD3,0xE8), False)]],
     line_spacing=1.12, space_after=6)

import os
_OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LoanWhiz-Presentation.pptx")
prs.save(_OUT)
print("saved", _OUT, "with", len(prs.slides._sldIdLst), "slides")
